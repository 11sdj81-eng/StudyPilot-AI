"""PPT/PPTX image extractor for StudyPilot Figure Engine.

Extracts images from PowerPoint files using python-pptx.  If python-pptx is not
installed, reports the issue gracefully instead of failing.
"""

from __future__ import annotations

import hashlib
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from core.figure_engine.figure_objects import FigureObject, SourceType

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False


class PptFigureExtractor:
    """Extract images from PPT/PPTX slide decks."""

    def __init__(
        self,
        output_dir: str | Path = "data/figure_bank/_extracted",
        available: bool = PPTX_AVAILABLE,
    ):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.available = available
        self.figures: list[FigureObject] = []
        self._figure_counter = 0

    def extract_from_file(
        self,
        pptx_path: str | Path,
        source_type: str = SourceType.PPT,
        concept_id: str | None = None,
    ) -> list[FigureObject]:
        """Extract images from a single PPTX file."""
        if not self.available:
            print("[PPTR] python-pptx is not installed. Skipping PPT extraction.")
            print("[PPTR] Install with: pip install python-pptx")
            return []

        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        prs = Presentation(str(pptx_path))
        figures: list[FigureObject] = []

        try:
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_title = _extract_slide_title(slide)
                slide_text = _extract_slide_text(slide)

                for shape_num, shape in enumerate(slide.shapes):
                    figure = self._extract_shape_image(
                        shape, shape_num, slide_num, slide_title, slide_text,
                        pptx_path, source_type, concept_id,
                    )
                    if figure:
                        figures.append(figure)
        except Exception as e:
            print(f"  [WARN] Failed to process {pptx_path.name}: {e}")

        self.figures.extend(figures)
        return figures

    def extract_from_directory(
        self,
        directory: str | Path,
    ) -> list[FigureObject]:
        """Extract images from all PPTX files in a directory (recursive)."""
        directory = Path(directory)
        figures: list[FigureObject] = []

        for pptx_file in sorted(directory.rglob("*.pptx")):
            try:
                figures.extend(self.extract_from_file(pptx_file))
            except Exception as e:
                print(f"  [WARN] Failed to extract from {pptx_file.name}: {e}")

        for ppt_file in sorted(directory.rglob("*.ppt")):
            print(f"  [INFO] Legacy .ppt format not supported for {ppt_file.name}; convert to .pptx first.")

        return figures

    def _extract_shape_image(
        self,
        shape: Any,
        shape_num: int,
        slide_num: int,
        slide_title: str,
        slide_text: str,
        pptx_path: Path,
        source_type: str,
        concept_id: str | None,
    ) -> FigureObject | None:
        """Extract image from a shape if it contains one."""
        try:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                img_bytes = image.blob
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"

                pptx_name = pptx_path.name
                figure_id = f"fig_{pptx_name[:20].replace(' ', '_')}_s{slide_num}_{shape_num}"
                out_path = self.output_dir / f"{figure_id}.{ext}"
                out_path.write_bytes(img_bytes)

                # Generate tags from slide title and text
                tags = _generate_tags(slide_title, slide_text)

                return FigureObject(
                    figure_id=figure_id,
                    concept_id=concept_id,
                    source_type=source_type,
                    source_file=pptx_name,
                    source_page=slide_num,
                    bbox=None,
                    image_path=str(out_path.resolve()),
                    caption=slide_title or None,
                    ocr_text=slide_text[:500] if slide_text else None,
                    tags=tags,
                    quality_score=0.0,
                    match_score=0.0,
                    usability_score=0.0,
                    width=None,
                    height=None,
                    aspect_ratio=None,
                    has_text_overlap_risk=False,
                    has_low_resolution_risk=False,
                    has_noise_risk=False,
                    metadata={
                        "slide_title": slide_title,
                        "slide_num": slide_num,
                        "shape_num": shape_num,
                        "ext": ext,
                        "content_type": image.content_type,
                    },
                )
        except Exception:
            pass
        return None

    @staticmethod
    def check_availability() -> bool:
        """Check if python-pptx is available."""
        return PPTX_AVAILABLE


def _extract_slide_title(slide: Any) -> str:
    """Extract title text from a slide."""
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    # Try first text shape as fallback
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            # First meaningful text line
            text = shape.text.strip().split("\n")[0]
            if len(text) >= 2:
                return text[:200]
            break
    return ""


def _extract_slide_text(slide: Any) -> str:
    """Extract all text from a slide."""
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                texts.append(text)
    return " ".join(texts)[:1000]


def _generate_tags(slide_title: str, slide_text: str) -> list[str]:
    """Generate keyword tags from slide title and text."""
    combined = f"{slide_title} {slide_text}".lower()
    tags: list[str] = []

    keyword_map = {
        "高斯": "gauss",
        "通量": "flux",
        "闭合面": "closed_surface",
        "电荷": "charge",
        "镜像": "image_method",
        "接地": "grounded",
        "导体": "conductor",
        "边界": "boundary",
        "介质": "dielectric",
        "法向": "normal",
        "切向": "tangential",
        "电位": "potential",
        "梯度": "gradient",
        "等位": "equipotential",
        "电场": "electric_field",
        "场强": "field_intensity",
        "场线": "field_lines",
        "能量": "energy",
        "电容": "capacitor",
        "电位移": "electric_displacement",
        "静电场": "electrostatic",
    }

    for cn_word, en_tag in keyword_map.items():
        if cn_word in combined:
            tags.append(en_tag)

    # Also keep original Chinese keywords
    for cn_word in keyword_map:
        if cn_word in combined:
            tags.append(cn_word)

    return list(set(tags))  # dedup
